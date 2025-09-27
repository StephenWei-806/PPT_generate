import json
import sys
import os
import json
from pptx import Presentation
from pptx.dml.color import RGBColor

# 动态获取资源文件路径（支持开发环境和PyInstaller打包环境）
def get_resource_path(relative_path):
    if hasattr(sys, '_MEIPASS'):
        # 打包后路径（资源文件位于_MEIPASS临时目录）
        return os.path.join(sys._MEIPASS, relative_path)
    # 开发环境路径（直接从当前工作目录获取）
    return os.path.join(os.path.abspath('.'), relative_path)

class Variable:
    REPEAT_PREFIX = '@repeat '
    
    @staticmethod
    def is_repeat(word):
        return word.startswith('{' + Variable.REPEAT_PREFIX)

class VarData:
    def __init__(self, data):
        self.data = data
        
    def get_var(self, var):
        if isinstance(var, str):
            var = var.strip('{}')
        return self.data.get(var)

def replace_text_in_ppt(json_data, ppt_filename):
    """替换PPT模板中的文本内容
    
    Args:
        json_data (str): JSON格式的PPT内容数据
        ppt_filename (str): PPT模板文件路径（绝对路径或相对路径）
        
    Returns:
        str: 生成的PPT文件路径
        
    Raises:
        FileNotFoundError: 当模板文件不存在时
        json.JSONDecodeError: 当JSON解析失败时
    """
    try:
        data = json.loads(json_data)
        print(f"成功解析JSON数据，包含键: {list(data.keys())}")
    except json.JSONDecodeError as e:
        print(f"JSON解析失败: {str(e)}")
        raise

    var_data = VarData(data)
    
    # 获取模板文件路径
    ppt_path = _resolve_template_path(ppt_filename)
    print(f"使用PPT模板路径: {ppt_path}")

    try:
        prs = Presentation(ppt_path)
        print(f"成功打开PPT模板，包含 {len(prs.slides)} 张幻灯片")
    except Exception as e:
        print(f"打开PPT模板失败: {str(e)}")
        print(f"尝试的模板路径: {ppt_path}")
        raise FileNotFoundError(f"无法打开PPT模板文件: {ppt_path}. 原因: {str(e)}")

    slides_to_delete = []

    def fill_paragraph(paragraph, data: VarData):
        def process(expr):
            value = data.get_var(expr)
            return str(value) if value is not None else ''

        var_name = ''
        text = ''
        for run in paragraph.runs:
            old_text = run.text
            run_text = ''
            for c in old_text:
                if not var_name:
                    if c == '{':
                        var_name = c
                    else:
                        run_text += c
                else:
                    var_name += c
                    if c == '}':
                        v = var_name.strip('{}')
                        if Variable.is_repeat(var_name):
                            key = v[len(Variable.REPEAT_PREFIX):]
                            replacement = data.get_var(key)
                            if replacement is None or replacement == "":
                                return None
                            run_text += str(replacement)
                        else:
                            if v in data.data:
                                print(f"替换变量 {{ {v} }} 为: {process(v)}")
                                run_text += process(v)
                            else:
                                run_text += var_name
                        var_name = ''
            
            if run_text:
                new_run = paragraph.add_run()
                new_run.text = run_text
                
                # 复制格式
                if hasattr(run.font, 'bold'):
                    new_run.font.bold = run.font.bold
                if hasattr(run.font, 'italic'):
                    new_run.font.italic = run.font.italic
                if hasattr(run.font, 'underline'):
                    new_run.font.underline = run.font.underline
                if hasattr(run.font, 'strike'):
                    new_run.font.strike = run.font.strike
                if hasattr(run.font, 'size'):
                    new_run.font.size = run.font.size
                if hasattr(run.font.color, 'rgb') and isinstance(run.font.color, RGBColor):
                    new_run.font.color.rgb = run.font.color.rgb
                
                text += run_text
            run.text = ''
        return paragraph

    # 处理幻灯片内容
    for slide_idx, slide in enumerate(prs.slides):
        delete_slide = False
        print(f"处理幻灯片 {slide_idx + 1}...")
        for shape_idx, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                print(f"  处理文本框 {shape_idx + 1}...")
                for paragraph_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    print(f"    处理段落 {paragraph_idx + 1}...")
                    result = fill_paragraph(paragraph, var_data)
                    if result is None:
                        delete_slide = True
                        print(f"    段落 {paragraph_idx + 1} 包含空的重复变量，标记幻灯片删除")
                        break
                if delete_slide:
                    break
        if delete_slide:
            slides_to_delete.append(slide_idx)
            print(f"幻灯片 {slide_idx + 1} 标记为删除")
    
    # 删除不需要的幻灯片
    for idx in sorted(slides_to_delete, reverse=True):
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        del slides[idx]
        xml_slides[:] = slides
        print(f"已删除幻灯片 {idx + 1}")
    
    # 保存生成的PPT到backend/temp/generated目录
    temp_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', 'temp', 'generated')
    # 确俜temp/generated目录存在
    os.makedirs(temp_dir, exist_ok=True)
    
    # 生成唯一的文件名
    import time
    timestamp = str(int(time.time()))
    new_filename = os.path.join(temp_dir, f'modified_{timestamp}.pptx')
    
    try:
        prs.save(new_filename)
        print(f"成功保存PPT文件到: {new_filename}")
    except Exception as e:
        print(f"保存PPT文件失败: {str(e)}")
        raise

    return new_filename

def _resolve_template_path(template_path):
    """解析模板文件路径
    
    Args:
        template_path (str): 模板文件路径（绝对路径或相对路径）
        
    Returns:
        str: 解析后的绝对路径
        
    Raises:
        FileNotFoundError: 当模板文件不存在时
    """
    # 如果是绝对路径且文件存在，直接返回
    if os.path.isabs(template_path) and os.path.exists(template_path):
        return template_path
    
    # 如果是相对路径，尝试多个位置
    search_paths = [
        # 直接使用原路径
        template_path,
        # 在当前工作目录中查找
        os.path.join(os.getcwd(), template_path),
        # 使用get_resource_path函数查找（兼容打包环境）
        get_resource_path(template_path),
        # 在backend目录中查找
        os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), template_path),
        # 在backend/temp/default目录中查找
        os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), 'temp', 'default', os.path.basename(template_path))
    ]
    
    for path in search_paths:
        if os.path.exists(path):
            print(f"[调试] 找到模板文件: {path}")
            return path
    
    # 所有路径都不存在
    raise FileNotFoundError(f"找不到模板文件: {template_path}。已尝试的路径: {search_paths}")